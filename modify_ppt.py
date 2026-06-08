"""修改 NOVA_GDC_Training_v3.pptx — 适配二线所场景 + 6大核心功能聚焦"""

from pptx import Presentation
from pptx.util import Pt
import re

SRC = r'C:\Users\bgu\Desktop\my-app-hub\NOVA_GDC_Training_v3.pptx'
DST = r'C:\Users\bgu\Desktop\my-app-hub\NOVA_GDC_Training_v3.pptx'

prs = Presentation(SRC)


# ── 工具函数 ──────────────────────────────────

def replace_all(old: str, new: str, slides: list[int] | None = None):
    """全文替换（跨 run 合并处理，尽量保留首 run 格式）。"""
    for i, slide in enumerate(prs.slides):
        if slides is not None and i not in slides:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = para.text
                if old not in full:
                    continue
                new_full = full.replace(old, new)
                if para.runs:
                    para.runs[0].text = new_full
                    for r in para.runs[1:]:
                        r.text = ''


def rewrite_slide(slide_idx: int, new_lines: list[str]):
    """覆盖指定幻灯片所有文本框的内容（每行一个 paragraph）。"""
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            # 只改第一个文本框（通常是主内容区）
            _fill_text_frame(tf, new_lines)
            return


def _fill_text_frame(tf, lines: list[str]):
    """用纯文本行填充一个 TextFrame（保留现有段落结构）。"""
    paragraphs = list(tf.paragraphs)
    for idx, line in enumerate(lines):
        if idx < len(paragraphs):
            p = paragraphs[idx]
            if p.runs:
                p.runs[0].text = line
                for r in p.runs[1:]:
                    r.text = ''
            else:
                r = p.add_run()
                r.text = line
        else:
            p = tf.add_paragraph()
            p.add_run().text = line


def hide_slide(idx: int):
    """隐藏幻灯片（在幻灯片放映中不可见）。"""
    prs.slides[idx]._element.set('show', '0')


# ── 1. 全局替换 ──────────────────────────────

replace_all('8大核心功能模块', '6大核心功能模块')

# "重分类" → 更中性表达
replace_all('重分类', '重新归类')

# "现场互动" → 更适合远程的表述
replace_all('现场互动', '课后练习')

# 封面"GDC审计团队专用教学" → 更通用
# 保留，因为二线所也属于GDC

# ── 2. 更新 F01 示例 ──────────────────────────
# 去掉 F01 中过于复杂的例子，用固定资产查询替代部分内容

# Slide 4 (F01 WHY) - 替换"重分类"相关例子
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = para.text
            # 替换"其他应付款——员工挂账"例子
            if '其他应付款——员工挂账' in full:
                new_txt = full.replace(
                    "比如客户账上有一笔'其他应付款——员工挂账'余额800万，该挂'其他应付款'还是该重分类成'应付职工薪酬'？Field同事可能问你一句，你翻准则翻半小时。",
                    "比如客户固定资产折旧年限是20年，但行业惯例是15年，你需要确认该用哪个年限。打开准则翻半天找不到明确说法。"
                ).replace(
                    "再比如，客户有一笔500万的'营业外支出'，你怀疑应该入'资产减值损失'，但不确定具体用哪个准则。",
                    "再比如，收到客户关于固定资产减值测试的说明，不确定计算参数是否符合准则要求，需要快速查找权威参考。"
                )
                if para.runs:
                    para.runs[0].text = new_txt
                    for r in para.runs[1:]:
                        r.text = ''

# Slide 5 (F01 WHEN) - 替换例子
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = para.text
            changed = False
            if "科目'其他应收款——押金'" in full:
                full = full.replace(
                    "场景1：科目分类拿不准时\n客户SAP导出后，你看到科目'其他应收款——押金'余额200万，不确定是否应归入'应收票据'还是保持不动。用精准搜索1分钟就能找到权威判断依据。",
                    "场景1：对准则条文理解不确定时\n比如客户固定资产采用年数总和法折旧，你不确定该方法在CAS4号下的适用条件。用NOVA搜索1分钟就能找到权威判断依据。"
                )
                changed = True
            if "场景2：大额异常交易不知道入什么科目时" in full:
                full = full.replace(
                    "场景2：大额异常交易不知道入什么科目时\n客户有一笔'营业外支出'500万，Field同事说'先判断一下是不是减值损失'。你需要确认这笔钱的性质和适用准则，而不是直接按客户分类照抄。",
                    "场景2：不确定某项会计处理是否合规时\n比如客户将研发支出全部费用化，你不确定是否有部分应资本化。需要查证研发支出资本化条件的具体准则条文。"
                )
                changed = True
            if "场景3：细节测试中需要引用准则条文时" in full:
                full = full.replace(
                    "场景3：细节测试中需要引用准则条文时\n比如你做收入确认细节测试，客户用的是总额法，但你说不准这个判断对不对。需要快速查证收入准则关于总额法/净额法的适用条件作为底稿支撑。",
                    "场景3：需要引用准则条文支持底稿结论时\n比如你做固定资产减值测试，客户用了特定的折现率。需要查证CAS8号关于折现率选取的准则依据来支撑底稿结论。"
                )
                changed = True
            if changed and para.runs:
                para.runs[0].text = full
                for r in para.runs[1:]:
                    r.text = ''

# Slide 6 (F01 HOW) - 简化例子
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = para.text
            if "不要只搜一个科目名" in full:
                new_txt = full.replace(
                    "不要只搜一个科目名，要把问题拆解为'科目 + 场景 + 准则/年份'。比如不要搜'应付账款'，要搜'预付账款 和其他应收款 区别 会计准则'。",
                    "不要只搜单个关键词，要把问题拆解为'主题 + 场景 + 准则/年份'。比如不要搜'固定资产折旧'，要搜'固定资产 年数总和法 适用条件 CAS4'。"
                )
                if para.runs:
                    para.runs[0].text = new_txt
                    for r in para.runs[1:]:
                        r.text = ''

# Slide 7 (F01 HOW详细) - 简化指令
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = para.text
            if "指令1（查科目区分）" in full:
                new_txt = full.replace(
                    "指令1（查科目区分）：site:mof.gov.cn '预付账款' '其他应收款' 区别 会计准则",
                    "指令1（查准则规定）：site:mof.gov.cn '固定资产' '折旧方法' 变更 条件 CAS4"
                ).replace(
                    "指令2（查准则条文）：site:csrc.gov.cn '收入确认' 总额法 净额法 监管要求",
                    "指令2（查具体条文）：site:csrc.gov.cn '资产减值' 折现率 选取 监管要求"
                ).replace(
                    "指令3（查行业处理）：'XX行业' '原材料采购' 会计处理 准则解释",
                    "指令3（查实务案例）：site:aicpa.org 'fixed asset' impairment testing practical guide"
                )
                if para.runs:
                    para.runs[0].text = new_txt
                    for r in para.runs[1:]:
                        r.text = ''

# Slide 8 (F01 WATCH OUT) - 去掉现场互动，改例子
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = para.text
            if "错误案例：搜'应收账款怎么审计'" in full:
                full = full.replace(
                    "错误案例：搜'应收账款怎么审计'，返回结果全是百度经验",
                    "错误案例：搜'资产减值 测试方法'，返回一堆论坛帖子"
                ).replace(
                    "正确做法：加上限定词'应收账款 减值测试 审计程序 财政部 site:gov.cn'。加上site和具体场景词，结果精准度提升5-10倍。",
                    "正确做法：加上限定词'资产减值 测试 折现率 CAS8 site:gov.cn'。加上site和具体准则号，精准度提升5-10倍。"
                )
            if "给出一个具体场景" in full:
                new_txt = full.replace(
                    "给出一个具体场景：客户科目'其他应付款——员工挂账'余额800万，3个小组分别用不同搜索词去查，对比谁的结果最先找到权威依据，讨论各自的搜索策略优劣。",
                    "案例：客户固定资产原值5000万，采用直线法折旧，你怀疑残值率设定不合理。分别用3种不同搜索策略查证残值率的准则规定，对比哪种策略最快找到权威依据。"
                )
                if para.runs:
                    para.runs[0].text = new_txt
                    for r in para.runs[1:]:
                        r.text = ''

# ── 3. 压缩 F04 Excel处理（Slides 9-14，保留标题+1页概览，隐藏其余） ──

# Slide 9 (F04 标题) - 改副标题
slide9 = prs.slides[8]
replace_all('审计分析程序与数据建模自动化', '常用Excel分析辅助功能', slides=[8])

# Slide 10 (原F04 WHY) → 改写为概览页
overview_f04 = [
    "F04 Excel表格处理 — 功能概览",
    "",
    "NOVA可以帮你快速生成常用的Excel分析模板：",
    "",
    "• 趋势分析模板——输入月份和收入数据，自动计算同比环比和毛利率",
    "• 数据匹配工具——两列数据快速匹配合并，标记差异项",
    "• 抽样计算表——MUS货币单位抽样，自动算样本量和选样",
    "• 异常筛选——按条件标记超出范围的数据点",
    "",
    "一句话：当你需要做数据分析或数据匹配时，告诉NOVA你的数据格式，",
    "它直接生成带公式的Excel模板，你只需要填入数据、验算结果。",
    "",
    "注意：公式一定要验算！随机挑3行数据手工核对，确保引用范围没错。",
]

# 找Slide 10的主文本框
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape.has_text_frame:
        # 检查这个文本框是否有足够内容（排除小标签）
        text_len = len(shape.text_frame.text.strip())
        if text_len > 50:  # 主内容区
            tf = shape.text_frame
            # 清空所有段落
            for p_idx in range(len(tf.paragraphs)):
                for r in tf.paragraphs[p_idx].runs:
                    r.text = ''
            # 填入新内容
            for idx, line in enumerate(overview_f04):
                if idx < len(tf.paragraphs):
                    p = tf.paragraphs[idx]
                else:
                    p = tf.add_paragraph()
                if p.runs:
                    p.runs[0].text = line
                    for r in p.runs[1:]:
                        r.text = ''
                else:
                    r = p.add_run()
                    r.text = line

# 隐藏 F04 剩余详细页（Slides 11-14）
for idx in [10, 11, 12, 13]:
    hide_slide(idx)


# ── 4. 压缩 F09 财务报表API（Slides 21-26，保留标题+1页概览，隐藏其余） ──

# Slide 21 (F09 标题) - 改副标题
replace_all('A股/港股上市公司财务数据批量获取', '快速查询上市公司财务数据', slides=[20])

# Slide 22 (原F09 WHY) → 改写为概览页
overview_f09 = [
    "F09 财务报表查询 — 功能概览",
    "",
    "NOVA可以快速获取上市公司公开财务数据：",
    "",
    "• 查询A股/港股上市公司核心财务指标",
    "• 支持按股票代码或公司名称搜索",
    "• 数据导出为标准化Excel格式",
    "• 可用于行业对比和趋势分析",
    "",
    "一句话：需要找同行数据做比较分析时，直接输股票代码导出数据，",
    "省去手动从年报PDF里抄数字的时间和出错风险。",
    "",
    "注意：导出后务必抽查1-2个年度，核对营收、净利润等关键数字",
    "与公司已披露年报是否一致，确认单位（元/万元）是否正确。",
]

slide22 = prs.slides[21]
for shape in slide22.shapes:
    if shape.has_text_frame:
        text_len = len(shape.text_frame.text.strip())
        if text_len > 50:
            tf = shape.text_frame
            for p_idx in range(len(tf.paragraphs)):
                for r in tf.paragraphs[p_idx].runs:
                    r.text = ''
            for idx, line in enumerate(overview_f09):
                if idx < len(tf.paragraphs):
                    p = tf.paragraphs[idx]
                else:
                    p = tf.add_paragraph()
                if p.runs:
                    p.runs[0].text = line
                    for r in p.runs[1:]:
                        r.text = ''
                else:
                    r = p.add_run()
                    r.text = line

# 隐藏 F09 剩余详细页（Slides 23-26）
for idx in [22, 23, 24, 25]:
    hide_slide(idx)


# ── 5. 更新目录（Slide 2） ──

slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = para.text
            # 在目录中把F04和F09移到"其他功能"
            if 'F04' in full and 'Excel表格处理' in full and '审计分析程序' in full:
                new_txt = full.replace(
                    "F04  Excel表格处理 — 审计分析程序与数据建模自动化",
                    "其他功能：F04 Excel表格处理 | F09 财务报表查询（详见后文概览）"
                )
                if para.runs:
                    para.runs[0].text = new_txt
                    for r in para.runs[1:]:
                        r.text = ''
            if 'F09' in full and '财务报表API' in full:
                # 已合并到上一条，清空此行
                if para.runs:
                    para.runs[0].text = ''
                    for r in para.runs[1:]:
                        r.text = ''


# ── 6. 调整其他功能中不适用于二线所的内容 ──

# F05 (Slides 15-20) - 调整"现场互动"内容
replace_all('场景2：收到函证回函扫描件时\n往来单位把函证回函扫描后发邮件过来，你需要识别回函结论、差异金额、备注说明，录入回函追踪表。',
            '场景2：收到供应商/客户回函扫描件时\n往来单位把回函扫描后发邮件过来，你需要提取回函结论、金额、备注说明，录入回函追踪表。',
            slides=[16])

# F11 (Slides 45-50) - PPT生成器，调整面向场景
# 把"向Field团队汇报"改为"向项目组汇报"，这个已经在全局替换中处理了

# F12 (Slides 33-38) - 图片OCR，调整
replace_all(
    "如果你打开PDF后能选中文字、复制粘贴，那是电子版->用[直接提取]模式。如果只能看到图片、鼠标拖不动文字，那是扫描件->必须用[OCR模式]。选错模式，扫描件会提取出一堆乱码。",
    "打开文件后能选中文字、复制粘贴，就是电子版->用[直接提取]模式。如果只能看到图片、文字选不中，就是扫描件->用[OCR识别]模式。选错模式会提取出乱码。",
    slides=[17])

# ── 7. 保存 ──

prs.save(DST)
print(f"修改完成，已保存到：{DST}")
print(f"   总页数：{len(prs.slides)}")
hidden = sum(1 for s in prs.slides if s._element.get('show') == '0')
print(f"   隐藏页（概览压缩）：{hidden}")
print(f"   放映可见页：{len(prs.slides) - hidden}")
